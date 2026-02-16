from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from pdf2docx import Converter
from deep_translator import GoogleTranslator
import concurrent.futures
import os

def translate_text(text, target_lang='or'):
    """Translates text to the target language using Google Translator."""
    if text is None:
        return ""
    if not text.strip():
        return text
    try:
        translated = GoogleTranslator(source='auto', target=target_lang).translate(text)
        if translated is None:
            return text
        return translated
    except Exception as e:
        print(f"Error translating text: {text[:20]}... Error: {e}")
        return text

def collect_runs(shape, runs_list):
    """
    Recursively processes a shape to collect text runs.
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for member in shape.shapes:
            collect_runs(member, runs_list)
        return

    # Handle text frames (includes placeholders)
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text and run.text.strip():
                    runs_list.append(run)
    
    # Handle tables
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                runs_list.append(run)
    
    # Handle charts
    if hasattr(shape, 'has_chart') and shape.has_chart:
        try:
            chart = shape.chart
            
            # Chart title
            if chart.has_title and chart.chart_title.has_text_frame:
                for paragraph in chart.chart_title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            runs_list.append(run)
            
            # Category axis title
            if hasattr(chart, 'category_axis') and chart.category_axis.has_title:
                axis_title = chart.category_axis.axis_title
                if hasattr(axis_title, 'has_text_frame') and axis_title.has_text_frame:
                    for paragraph in axis_title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                runs_list.append(run)
            
            # Value axis title
            if hasattr(chart, 'value_axis') and chart.value_axis.has_title:
                axis_title = chart.value_axis.axis_title
                if hasattr(axis_title, 'has_text_frame') and axis_title.has_text_frame:
                    for paragraph in axis_title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text and run.text.strip():
                                runs_list.append(run)
            
            # Legend (if present)
            if chart.has_legend:
                # Note: Legend entries are typically auto-generated from data
                # and may not have editable text frames in all chart types
                pass
                
        except Exception as e:
            print(f"Error processing chart: {e}")

def translate_pptx(input_path, output_path, target_lang='or'):
    """
    Parses a PPTX file, translates its text content to the target language using threading,
    and saves the translated file.
    """
    try:
        prs = Presentation(input_path)
    except Exception as e:
        print(f"Error loading presentation: {e}")
        raise e

    # 1. Collect all runs with text from slides, master slides, and layouts
    all_runs = []
    
    # Process regular slides
    for slide in prs.slides:
        for shape in slide.shapes:
            collect_runs(shape, all_runs)
    
    # Process slide masters and their layouts
    for slide_master in prs.slide_masters:
        # Process shapes in the master slide itself
        for shape in slide_master.shapes:
            collect_runs(shape, all_runs)
        
        # Process each layout in the master
        for layout in slide_master.slide_layouts:
            for shape in layout.shapes:
                collect_runs(shape, all_runs)

    print(f"Found {len(all_runs)} text runs to translate.")

    # 2. Extract unique texts to translate (optimization)
    unique_texts = list(set([run.text for run in all_runs]))
    print(f"Found {len(unique_texts)} unique text strings.")

    # 3. Translate unique texts in parallel
    translations = {}
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        future_to_text = {executor.submit(translate_text, text, target_lang): text for text in unique_texts}
        for future in concurrent.futures.as_completed(future_to_text):
            original_text = future_to_text[future]
            try:
                translated_text = future.result()
                translations[original_text] = translated_text
            except Exception as e:
                print(f"Error translating {original_text}: {e}")
                translations[original_text] = original_text

    # 4. Apply translations back to runs
    for run in all_runs:
        if run.text in translations:
            run.text = translations[run.text]

    try:
        prs.save(output_path)
        print(f"Translated presentation saved to {output_path}")
        return output_path
    except Exception as e:
        print(f"Error saving presentation: {e}")
        raise e

def translate_docx(input_path, output_path, target_lang='or'):
    """
    Parses a DOCX file, translates its text content, and saves it.
    """
    try:
        doc = Document(input_path)
    except Exception as e:
        print(f"Error loading document: {e}")
        raise e

    # Collect all text to translate from paragraphs and tables
    all_texts = []
    
    # Text from paragraphs
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if run.text and run.text.strip():
                all_texts.append(run)

    # Text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if run.text and run.text.strip():
                            all_texts.append(run)

    print(f"Found {len(all_texts)} text runs to translate in DOCX.")

    # Extract unique texts
    unique_texts = list(set([run.text for run in all_texts]))
    print(f"Found {len(unique_texts)} unique text strings in DOCX.")

    # Translate unique texts in parallel
    translations = {}
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        future_to_text = {executor.submit(translate_text, text, target_lang): text for text in unique_texts}
        for future in concurrent.futures.as_completed(future_to_text):
            original_text = future_to_text[future]
            try:
                translated_text = future.result()
                translations[original_text] = translated_text
            except Exception as e:
                print(f"Error translating {original_text}: {e}")
                translations[original_text] = original_text

    # Apply translations
    for run in all_texts:
        if run.text in translations:
            run.text = translations[run.text]

    try:
        doc.save(output_path)
        print(f"Translated document saved to {output_path}")
        return output_path
    except Exception as e:
        print(f"Error saving document: {e}")
        raise e

def translate_pdf(input_path, output_path, target_lang='or'):
    """
    Converts PDF to DOCX, translates the DOCX, and saves it.
    The output is a DOCX file.
    """
    # 1. Convert PDF to DOCX
    docx_temp_path = input_path.replace('.pdf', '_temp.docx')
    try:
        cv = Converter(input_path)
        cv.convert(docx_temp_path, start=0, end=None)
        cv.close()
        print(f"Converted PDF to DOCX: {docx_temp_path}")
    except Exception as e:
        print(f"Error converting PDF to DOCX: {e}")
        raise e

    # 2. Translate the converted DOCX
    # output_path is essentially the final translated DOCX
    try:
        translate_docx(docx_temp_path, output_path, target_lang)
    except Exception as e:
        print(f"Error translating converted DOCX: {e}")
        raise e
    finally:
        # Cleanup temp file
        if os.path.exists(docx_temp_path):
            os.remove(docx_temp_path)

    return output_path

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 2:
        input_file = sys.argv[1]
        output_file = sys.argv[2]
        
        if input_file.lower().endswith('.pptx'):
            translate_pptx(input_file, output_file)
        elif input_file.lower().endswith('.docx'):
            translate_docx(input_file, output_file)
        elif input_file.lower().endswith('.pdf'):
            translate_pdf(input_file, output_file)
        else:
            print("Unsupported file format for CLI usage.")
    else:
        print("Usage: python translator.py <input_file> <output_file>")
