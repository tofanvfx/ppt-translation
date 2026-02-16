from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def create_test_pptx(filename):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    # Create a group
    shapes = slide.shapes
    group_shape = shapes.add_group_shape()
    
    # Add shapes to group
    # Note: python-pptx add_group_shape returns a group, but adding shapes to it is different in newer versions or might require specific methods.
    # Actually, getting a group shape in python-pptx and adding to it programmatically is tricky/not fully supported in all versions.
    # Instead, let's just create a shape and see if we can check properties.
    
    # We can't easily create a group programmatically with robust content in older python-pptx versions without complex xml manipulation.
    # But we can check if MSO_SHAPE_TYPE is importable.
    print("MSO_SHAPE_TYPE imported successfully")
    return filename

if __name__ == "__main__":
    try:
        create_test_pptx("test.pptx")
        print("Success")
    except Exception as e:
        print(f"Error: {e}")
